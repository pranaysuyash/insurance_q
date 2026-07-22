"""Tests for source-preserving native DOCX observations."""

from __future__ import annotations

import io

import pytest
from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches as PptxInches

from src.ocr.native_office import (
    extract_native_docx_document,
    extract_native_pptx_document,
    extract_native_xlsx_document,
)
from src.services.document_processing_service import DocumentProcessingService


def _docx_fixture() -> bytes:
    document = Document()
    document.add_heading("Health Policy", level=1)
    document.add_paragraph("Coverage is subject to the policy schedule.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Benefit"
    table.cell(0, 1).text = "Limit"
    table.cell(1, 0).text = "Hospitalization"
    table.cell(1, 1).text = "500000"

    image = Image.new("RGB", (16, 16), "#2f80ed")
    image_buffer = io.BytesIO()
    image.save(image_buffer, format="PNG")
    document.add_picture(io.BytesIO(image_buffer.getvalue()), width=Inches(0.2))

    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _xlsx_fixture() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Benefits"
    sheet.append(["Benefit", "Limit"])
    sheet.append(["Hospitalization", 500000])
    sheet.append(["Formula", "=B2*2"])
    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def _pptx_fixture() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Health Policy"
    textbox = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(5), PptxInches(1))
    textbox.text = "Coverage is subject to the policy schedule."
    table_shape = slide.shapes.add_table(2, 2, PptxInches(1), PptxInches(3), PptxInches(5), PptxInches(1))
    table_shape.table.cell(0, 0).text = "Benefit"
    table_shape.table.cell(0, 1).text = "Limit"
    table_shape.table.cell(1, 0).text = "Hospitalization"
    table_shape.table.cell(1, 1).text = "500000"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_native_docx_adapter_preserves_structure_and_image_lineage():
    result = extract_native_docx_document(_docx_fixture())

    assert result["parser_profile"] == "python-docx-native"
    assert "Health Policy" in result["full_text"]
    assert "500000" in result["full_text"]
    node_types = {node.node_type for node in result["nodes"]}
    assert {"page", "heading", "text_block", "table", "table_cell", "figure"} <= node_types

    heading = next(node for node in result["nodes"] if node.node_type == "heading")
    assert heading.attributes["heading_level"] == 1
    cell = next(node for node in result["nodes"] if node.node_type == "table_cell")
    assert cell.attributes["row_index"] == 0
    assert cell.attributes["column_index"] == 0
    figure = next(node for node in result["nodes"] if node.node_type == "figure")
    assert len(figure.artifact_sha256) == 64


@pytest.mark.asyncio
async def test_document_service_uses_native_docx_path(tmp_path):
    source = tmp_path / "policy.docx"
    source.write_bytes(_docx_fixture())

    service = DocumentProcessingService()
    result = await service._extract_text(str(source), "policy.docx")

    assert result["status"] == "completed"
    assert result["method"] == "python-docx-native"
    assert result["cir"]["parser_profile"] == "python-docx-native"
    assert any(node["node_type"] == "heading" for node in result["cir"]["nodes"])


def test_native_xlsx_adapter_preserves_cells_and_formulas():
    result = extract_native_xlsx_document(_xlsx_fixture())

    assert result["parser_profile"] == "openpyxl-native"
    assert "=B2*2" in result["full_text"]
    cells = [node for node in result["nodes"] if node.node_type == "table_cell"]
    formula = next(node for node in cells if node.source_text == "=B2*2")
    assert formula.attributes["coordinate"] == "B3"
    assert formula.attributes["is_formula"] is True


def test_native_pptx_adapter_preserves_slide_structure():
    result = extract_native_pptx_document(_pptx_fixture())

    assert result["parser_profile"] == "python-pptx-native"
    assert "Health Policy" in result["full_text"]
    assert "500000" in result["full_text"]
    node_types = {node.node_type for node in result["nodes"]}
    assert {"page", "heading", "text_block", "table", "table_cell"} <= node_types


@pytest.mark.asyncio
async def test_document_service_uses_native_xlsx_and_pptx_paths(tmp_path):
    service = DocumentProcessingService()
    for filename, content, method in (
        ("policy.xlsx", _xlsx_fixture(), "openpyxl-native"),
        ("policy.pptx", _pptx_fixture(), "python-pptx-native"),
    ):
        source = tmp_path / filename
        source.write_bytes(content)
        result = await service._extract_text(str(source), filename)
        assert result["status"] == "completed"
        assert result["method"] == method
        assert result["cir"]["parser_profile"] == method
