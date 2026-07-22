"""Source-preserving native Office observations.

These adapters preserve native Office structure into the existing CIR; they do
not OCR images or infer semantic insurance fields from labels.
"""

from __future__ import annotations

import hashlib
import io
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.models.document_intelligence import CIRNode


def _node(
    *,
    node_id: str,
    node_type: str,
    source_text: str | None = None,
    parent_id: str | None = None,
    reading_order: int,
    attributes: dict[str, Any] | None = None,
    artifact_sha256: str | None = None,
    page_number: int = 1,
    parser_name: str = "python-docx-native",
) -> CIRNode:
    return CIRNode(
        node_id=node_id,
        node_type=node_type,
        page_number=page_number,
        source_text=source_text,
        retrieval_text=source_text,
        parent_id=parent_id,
        reading_order=reading_order,
        artifact_sha256=artifact_sha256,
        parser_name=parser_name,
        parser_version=f"{parser_name}.v1",
        attributes=attributes or {},
    )


def extract_native_docx_document(docx_bytes: bytes) -> dict[str, Any]:
    """Extract DOCX structure without flattening it into binary text."""

    source_hash = hashlib.sha256(docx_bytes).hexdigest()
    document = Document(io.BytesIO(docx_bytes))
    nodes: list[CIRNode] = [
        _node(
            node_id="page-1",
            node_type="page",
            reading_order=0,
            artifact_sha256=source_hash,
            attributes={"source": "docx_document", "page_model": "logical_document"},
        )
    ]
    text_parts: list[str] = []
    order = 1

    for paragraph_index, paragraph in enumerate(document.paragraphs):
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = getattr(paragraph.style, "name", "") or ""
        is_heading = style_name.lower().startswith("heading")
        node_type = "heading" if is_heading else "text_block"
        node_id = f"page-1-{node_type}-{paragraph_index}"
        nodes.append(
            _node(
                node_id=node_id,
                node_type=node_type,
                source_text=text,
                parent_id="page-1",
                reading_order=order,
                artifact_sha256=source_hash,
                attributes={
                    "style_name": style_name,
                    "source": "docx_paragraph",
                    "heading_level": (
                        int(style_name.split()[-1])
                        if is_heading and style_name.split()[-1].isdigit()
                        else None
                    ),
                },
            )
        )
        text_parts.append(text)
        order += 1

    for table_index, table in enumerate(document.tables):
        table_id = f"page-1-table-{table_index}"
        rows = len(table.rows)
        columns = max((len(row.cells) for row in table.rows), default=0)
        nodes.append(
            _node(
                node_id=table_id,
                node_type="table",
                parent_id="page-1",
                reading_order=order,
                artifact_sha256=source_hash,
                attributes={
                    "rows": rows,
                    "columns": columns,
                    "source": "docx_table",
                },
            )
        )
        order += 1
        for row_index, row in enumerate(table.rows):
            for column_index, cell in enumerate(row.cells):
                text = cell.text.strip()
                if text:
                    text_parts.append(text)
                nodes.append(
                    _node(
                        node_id=f"{table_id}-cell-{row_index}-{column_index}",
                        node_type="table_cell",
                        source_text=text or None,
                        parent_id=table_id,
                        reading_order=order,
                        artifact_sha256=source_hash,
                        attributes={
                            "table_id": table_id,
                            "row_index": row_index,
                            "column_index": column_index,
                            "source": "docx_table_cell",
                        },
                    )
                )
                order += 1

    image_index = 0
    for relationship in document.part.rels.values():
        if "image" not in relationship.reltype:
            continue
        image_bytes = relationship.target_part.blob
        if not image_bytes:
            continue
        nodes.append(
            _node(
                node_id=f"page-1-figure-{image_index}",
                node_type="figure",
                parent_id="page-1",
                reading_order=order,
                artifact_sha256=hashlib.sha256(image_bytes).hexdigest(),
                attributes={
                    "source": "docx_relationship",
                    "relationship_id": relationship.rId,
                    "mime_type": relationship.target_part.content_type,
                },
            )
        )
        image_index += 1
        order += 1

    return {
        "full_text": "\n\n".join(text_parts),
        "page_texts": {1: "\n\n".join(text_parts)},
        "nodes": nodes,
        "source_artifact_sha256": source_hash,
        "parser_profile": "python-docx-native",
    }


def _native_page_document(
    *,
    source_bytes: bytes,
    pages: list[tuple[str, str, list[CIRNode], str]],
    parser_profile: str,
) -> dict[str, Any]:
    """Assemble a multi-sheet/slide native document into the CIR contract."""

    source_hash = hashlib.sha256(source_bytes).hexdigest()
    nodes: list[CIRNode] = []
    page_texts: dict[int, str] = {}
    for page_number, (page_id, title, page_nodes, page_text) in enumerate(pages, start=1):
        nodes.append(
            _node(
                node_id=page_id,
                node_type="page",
                reading_order=0,
                artifact_sha256=source_hash,
                attributes={"source": parser_profile, "title": title},
                page_number=page_number,
                parser_name=parser_profile,
            )
        )
        nodes.extend(page_nodes)
        page_texts[page_number] = page_text
    return {
        "full_text": "\n\n".join(text for text in page_texts.values() if text),
        "page_texts": page_texts,
        "nodes": nodes,
        "source_artifact_sha256": source_hash,
        "parser_profile": parser_profile,
    }


def extract_native_xlsx_document(xlsx_bytes: bytes) -> dict[str, Any]:
    """Extract worksheet cells, formulas, and embedded image lineage."""

    workbook = load_workbook(io.BytesIO(xlsx_bytes), data_only=False, read_only=False)
    pages: list[tuple[str, str, list[CIRNode], str]] = []
    for sheet_number, worksheet in enumerate(workbook.worksheets, start=1):
        page_id = f"sheet-{sheet_number}"
        nodes: list[CIRNode] = []
        order = 1
        text_rows: list[str] = []
        min_row, min_column, max_row, max_column = (
            worksheet.min_row,
            worksheet.min_column,
            worksheet.max_row,
            worksheet.max_column,
        )
        if worksheet.max_row and worksheet.max_column:
            table_id = f"{page_id}-table-0"
            nodes.append(
                _node(
                    node_id=table_id,
                    node_type="table",
                    reading_order=order,
                    parent_id=page_id,
                    attributes={
                        "source": "xlsx_worksheet",
                        "sheet_name": worksheet.title,
                        "rows": max_row - min_row + 1,
                        "columns": max_column - min_column + 1,
                    },
                    page_number=sheet_number,
                    parser_name="openpyxl-native",
                )
            )
            order += 1
            for row in worksheet.iter_rows(
                min_row=min_row,
                max_row=max_row,
                min_col=min_column,
                max_col=max_column,
            ):
                values = [cell.value for cell in row]
                if any(value is not None for value in values):
                    text_rows.append(" | ".join(str(value) for value in values if value is not None))
                for cell in row:
                    value = cell.value
                    nodes.append(
                        _node(
                            node_id=f"{table_id}-cell-{cell.row}-{cell.column}",
                            node_type="table_cell",
                            source_text=None if value is None else str(value),
                            reading_order=order,
                            parent_id=table_id,
                            attributes={
                                "source": "xlsx_cell",
                                "sheet_name": worksheet.title,
                                "coordinate": cell.coordinate,
                                "row_index": cell.row,
                                "column_index": cell.column,
                                "data_type": cell.data_type,
                                "is_formula": isinstance(value, str) and value.startswith("="),
                            },
                            page_number=sheet_number,
                            parser_name="openpyxl-native",
                        )
                    )
                    order += 1
        for image_index, image in enumerate(getattr(worksheet, "_images", [])):
            image_bytes = image._data()
            anchor = getattr(image, "anchor", None)
            anchor_ref = getattr(anchor, "_from", None)
            nodes.append(
                _node(
                    node_id=f"{page_id}-figure-{image_index}",
                    node_type="figure",
                    reading_order=order,
                    parent_id=page_id,
                    artifact_sha256=hashlib.sha256(image_bytes).hexdigest(),
                    attributes={
                        "source": "xlsx_embedded_image",
                        "sheet_name": worksheet.title,
                        "anchor_row": getattr(anchor_ref, "row", None),
                        "anchor_column": getattr(anchor_ref, "col", None),
                    },
                    page_number=sheet_number,
                    parser_name="openpyxl-native",
                )
            )
            order += 1
        pages.append((page_id, worksheet.title, nodes, "\n".join(text_rows)))
    return _native_page_document(
        source_bytes=xlsx_bytes, pages=pages, parser_profile="openpyxl-native"
    )


def extract_native_pptx_document(pptx_bytes: bytes) -> dict[str, Any]:
    """Extract slide text, tables, and embedded picture lineage."""

    presentation = Presentation(io.BytesIO(pptx_bytes))
    pages: list[tuple[str, str, list[CIRNode], str]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        page_id = f"slide-{slide_number}"
        nodes: list[CIRNode] = []
        text_parts: list[str] = []
        order = 1
        title = f"Slide {slide_number}"
        for shape_index, shape in enumerate(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    is_title = getattr(shape, "is_placeholder", False) and getattr(
                        getattr(shape, "placeholder_format", None), "type", None
                    ) == 1
                    if is_title:
                        title = text
                    nodes.append(
                        _node(
                            node_id=f"{page_id}-text-{shape_index}",
                            node_type="heading" if is_title else "text_block",
                            source_text=text,
                            reading_order=order,
                            parent_id=page_id,
                            attributes={"source": "pptx_text_shape", "shape_index": shape_index},
                            page_number=slide_number,
                            parser_name="python-pptx-native",
                        )
                    )
                    text_parts.append(text)
                    order += 1
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_id = f"{page_id}-table-{shape_index}"
                nodes.append(
                    _node(
                        node_id=table_id,
                        node_type="table",
                        reading_order=order,
                        parent_id=page_id,
                        attributes={
                            "source": "pptx_table",
                            "rows": len(shape.table.rows),
                            "columns": len(shape.table.columns),
                        },
                        page_number=slide_number,
                        parser_name="python-pptx-native",
                    )
                )
                order += 1
                for row_index, row in enumerate(shape.table.rows):
                    for column_index, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        if text:
                            text_parts.append(text)
                        nodes.append(
                            _node(
                                node_id=f"{table_id}-cell-{row_index}-{column_index}",
                                node_type="table_cell",
                                source_text=text or None,
                                reading_order=order,
                                parent_id=table_id,
                                attributes={
                                    "source": "pptx_table_cell",
                                    "row_index": row_index,
                                    "column_index": column_index,
                                },
                                page_number=slide_number,
                                parser_name="python-pptx-native",
                            )
                        )
                        order += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_bytes = shape.image.blob
                nodes.append(
                    _node(
                        node_id=f"{page_id}-figure-{shape_index}",
                        node_type="figure",
                        reading_order=order,
                        parent_id=page_id,
                        artifact_sha256=hashlib.sha256(image_bytes).hexdigest(),
                        attributes={
                            "source": "pptx_picture",
                            "shape_index": shape_index,
                            "mime_type": shape.image.content_type,
                        },
                        page_number=slide_number,
                        parser_name="python-pptx-native",
                    )
                )
                order += 1
        pages.append((page_id, title, nodes, "\n".join(text_parts)))
    return _native_page_document(
        source_bytes=pptx_bytes, pages=pages, parser_profile="python-pptx-native"
    )
